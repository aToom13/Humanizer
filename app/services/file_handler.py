import os
import io
from docx import Document
from pptx import Presentation

class FileHandler:
    @staticmethod
    def allowed_file(filename):
        return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in {'txt', 'md', 'docx', 'pptx'}

    @staticmethod
    def parse_file(file_storage):
        filename = file_storage.filename
        ext = filename.rsplit('.', 1)[1].lower()
        content = ""

        try:
            if ext in ['txt', 'md']:
                content = file_storage.read().decode('utf-8', errors='ignore')
            
            elif ext == 'docx':
                doc = Document(file_storage)
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                content = '\n'.join(full_text)
            
            elif ext == 'pptx':
                prs = Presentation(file_storage)
                full_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            full_text.append(shape.text)
                content = '\n'.join(full_text)
                
        except Exception as e:
            raise Exception(f"Error parsing file: {str(e)}")
            
        return content

    @staticmethod
    def create_docx(text):
        doc = Document()
        # Add basic formatting
        for paragraph in text.split('\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph)
        
        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio

    @staticmethod
    def create_txt(text):
        return io.BytesIO(text.encode('utf-8'))
